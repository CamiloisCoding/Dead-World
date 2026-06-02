# -*- coding: utf-8 -*-
"""Erstellt 'Dead World.pptx' im Terminal/Amber-Stil des Spiels."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Farbpalette: hell & beamer-tauglich (hoher Kontrast) ----
BG          = RGBColor(0xF5, 0xF1, 0xE8)   # warmes Papierweiß (Folien-Hintergrund)
BG_PANEL    = RGBColor(0xFF, 0xFF, 0xFF)   # weiße Karten
AMBER       = RGBColor(0x2B, 0x25, 0x20)   # Haupttext (dunkel, gut lesbar)
AMBER_BRT   = RGBColor(0xA8, 0x16, 0x16)   # Überschriften / Akzent (kräftiges Rot)
AMBER_DIM   = RGBColor(0x8A, 0x7C, 0x66)   # Sekundär / Rahmen / Fußzeile (Taupe)
BLOOD       = RGBColor(0xA8, 0x16, 0x16)   # Akzentbalken
DANGER      = RGBColor(0xC0, 0x2A, 0x12)   # Gefahr / Rot
GREEN_OK    = RGBColor(0x1E, 0x7A, 0x34)   # Prompt-Marker / Erfolg (dunkelgrün)
WHITE       = RGBColor(0x2B, 0x25, 0x20)   # auf hellem BG -> dunkel
# Dunkles Kopfband (hoher Kontrast auf Beamer)
BAND        = RGBColor(0x1E, 0x18, 0x12)   # dunkles Band
BAND_TXT    = RGBColor(0xF2, 0xC9, 0x6A)   # Amber-Titel auf dunklem Band
BAND_PROMPT = RGBColor(0x7C, 0xD6, 0x8A)   # heller grüner Prompt auf Band

MONO = "Consolas"

prs = Presentation()
prs.slide_width  = Emu(12192000)   # 16:9
prs.slide_height = Emu(6858000)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_bg(slide, color=BG):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    slide.shapes._spTree.remove(s._element)
    slide.shapes._spTree.insert(2, s._element)
    return s


def add_box(slide, l, t, w, h, fill=None, line=None, line_w=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def txt(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (text,size,color,bold)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (text, size, color, bold) in para:
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = MONO
    return tb


def header(slide, prompt, title, idx=None):
    """Dunkles Terminal-Kopfband (kontraststark fuer Beamer)."""
    add_box(slide, 0, 0, SW, Inches(1.15), fill=BAND)
    add_box(slide, 0, Inches(1.15), SW, Pt(3), fill=BLOOD)
    txt(slide, Inches(0.55), Inches(0.16), Inches(11), Inches(0.85), [
        [(prompt, 16, BAND_PROMPT, True), (title, 30, BAND_TXT, True)],
    ], anchor=MSO_ANCHOR.MIDDLE)
    if idx is not None:
        txt(slide, SW - Inches(2.0), Inches(0.16), Inches(1.5), Inches(0.85),
            [[(idx, 14, AMBER_DIM, False)]], align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, n):
    txt(slide, Inches(0.55), SH - Inches(0.5), Inches(6), Inches(0.4),
        [[("DEAD WORLD  //  v1.0.0", 11, AMBER_DIM, False)]],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, SW - Inches(2.0), SH - Inches(0.5), Inches(1.45), Inches(0.4),
        [[(f"{n:02d}", 11, AMBER_DIM, False)]], align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE)


def bullet_para(text, size=18, color=AMBER, bold=False, marker="> "):
    return [(marker, size, GREEN_OK, True), (text, size, color, bold)]


# ============================================================
# SLIDE 1 — Titel
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
# Rahmen
add_box(s, Inches(0.35), Inches(0.35), SW - Inches(0.7), SH - Inches(0.7),
        line=AMBER_DIM, line_w=1.5)
# Blutbalken links
add_box(s, Inches(0.35), Inches(0.35), Pt(6), SH - Inches(0.7), fill=BLOOD)
txt(s, 0, Inches(0.7), SW, Inches(0.6),
    [[("[ SYSTEM REBOOT ... QUARANTÄNE-PROTOKOLL AKTIV ]", 15, AMBER_DIM, False)]],
    align=PP_ALIGN.CENTER)
txt(s, 0, Inches(2.25), SW, Inches(1.7),
    [[("DEAD  WORLD", 96, AMBER_BRT, True)]], align=PP_ALIGN.CENTER,
    anchor=MSO_ANCHOR.MIDDLE)
add_box(s, SW/2 - Inches(3), Inches(4.05), Inches(6), Pt(2), fill=BLOOD)
txt(s, 0, Inches(4.25), SW, Inches(0.8),
    [[("Ein Text-Survival-Horror im Terminal-Stil", 24, AMBER, False)]],
    align=PP_ALIGN.CENTER)
txt(s, 0, Inches(5.55), SW, Inches(0.6),
    [[("Ein Projekt von Camilo", 20, WHITE, True)]], align=PP_ALIGN.CENTER)
txt(s, 0, SH - Inches(1.0), SW, Inches(0.5),
    [[("Python  ·  Pygame  ·  ~9.000 Zeilen Code", 14, AMBER_DIM, False)]],
    align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2 — Gliederung
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "> ", "Gliederung")
points = [
    "Vorstellung des Spiels",
    "Spielidee & Konzept",
    "Ziele des Spielers",
    "Besondere Spielkonzepte",
    "NPCs & Begleiter",
    "Steuerung / Gameplay",
    "Technischer Aufbau",
    "Schwierigkeiten beim Entwickeln",
    "Was ich gelernt habe",
    "Zukunft / Verbesserungen",
    "Demo des Spiels",
]
runs = []
for i, p in enumerate(points, 1):
    runs.append([(f"{i:02d}  ", 18, AMBER_DIM, True), (p, 20, AMBER, False)])
txt(s, Inches(1.2), Inches(1.55), Inches(9.8), Inches(5.0), runs,
    space_after=9, line_spacing=1.0)
footer(s, 2)

# ============================================================
# SLIDE 3 — Vorstellung des Spiels
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "> ", "Vorstellung des Spiels")
# linke Spalte: Genre / Plattform
add_box(s, Inches(0.55), Inches(1.55), Inches(5.0), Inches(4.6), fill=BG_PANEL,
        line=AMBER_DIM)
txt(s, Inches(0.85), Inches(1.8), Inches(4.5), Inches(4.2), [
    [("GENRE", 20, AMBER_BRT, True)],
    bullet_para("Survival-Horror", 18),
    bullet_para("Text-Adventure (Parser)", 18),
    bullet_para("Post-Apokalypse / Zombies", 18),
    [("", 8, AMBER, False)],
    [("PLATTFORM", 20, AMBER_BRT, True)],
    bullet_para("PC (Windows)", 18),
    bullet_para("Python 3 + Pygame", 18),
    [("", 8, AMBER, False)],
    [("SPRACHE", 20, AMBER_BRT, True)],
    bullet_para("Komplett auf Deutsch", 18),
], space_after=7)
# rechte Spalte: Spielidee
add_box(s, Inches(5.85), Inches(1.55), Inches(5.75), Inches(4.6), fill=BG_PANEL,
        line=AMBER_DIM)
txt(s, Inches(6.15), Inches(1.8), Inches(5.2), Inches(4.2), [
    [("SPIELIDEE", 20, AMBER_BRT, True)],
    [("Eine Toxoplasma-Seuche hat die Welt in Zombies "
      "verwandelt. Du bist Albert Wesker Cristal und "
      "fliehst in einen Bunker.", 17, AMBER, False)],
    [("", 6, AMBER, False)],
    [("Über ein retro Amber-Terminal erkundest du per "
      "Texteingabe eine zerstörte Stadt: Krankenhaus, "
      "Bibliothek, Walddorf, Casino, ein Geheimlabor "
      "und mehr.", 17, AMBER, False)],
    [("", 6, AMBER, False)],
    [("Überleben, kämpfen, looten und das Geheimnis der "
      "Seuche aufdecken.", 17, AMBER_BRT, True)],
], space_after=6, line_spacing=1.05)
footer(s, 3)

# ============================================================
# SLIDE 4 — Spielidee & Konzept
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "> ", "Spielidee & Konzept")
txt(s, Inches(0.85), Inches(1.5), Inches(10.5), Inches(0.8), [
    [("Klassisches Text-Adventure trifft auf moderne Zombie-Apokalypse.",
      19, AMBER_BRT, True)],
])
cards = [
    ("ERKUNDEN", "~100 handgebaute Räume mit Himmelsrichtungen & Karte"),
    ("ÜBERLEBEN", "Lebenspunkte, Essen & Heilung managen"),
    ("KÄMPFEN", "Nah- & Fernkampf gegen Zombies und Bosse"),
    ("LOOTEN", "Waffen, Munition, Medkits & Container durchsuchen"),
    ("ATMOSPHÄRE", "Dynamische Musik, Zombie-Sounds, CRT-Terminal-Look"),
    ("GEHEIMNISSE", "Geheimlabor, versteckte Orte & eine düstere Story"),
]
cx, cy = Inches(0.55), Inches(2.45)
cw, ch = Inches(3.6), Inches(1.75)
gx, gy = Inches(0.18), Inches(0.2)
for i, (t, d) in enumerate(cards):
    col, row = i % 3, i // 3
    l = cx + col * (cw + gx)
    tp = cy + row * (ch + gy)
    add_box(s, l, tp, cw, ch, fill=BG_PANEL, line=AMBER_DIM)
    add_box(s, l, tp, Pt(5), ch, fill=BLOOD)
    txt(s, l + Inches(0.25), tp + Inches(0.15), cw - Inches(0.45), ch - Inches(0.3), [
        [(t, 18, AMBER_BRT, True)],
        [(d, 14, AMBER, False)],
    ], space_after=6, line_spacing=1.05)
footer(s, 4)

# ============================================================
# SLIDE 5 — Ziele des Spielers
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "> ", "Ziele des Spielers")
goals = [
    ("Überleben", "So lange wie möglich am Leben bleiben — Gesundheit und Vorräte im Blick behalten."),
    ("Die Stadt erkunden", "Sich durch ~100 Räume kämpfen und neue Gebiete freischalten."),
    ("Ausrüstung sammeln", "Bessere Waffen, Munition und Heilmittel finden."),
    ("Zombies besiegen", "Gegner ausschalten und Punkte (Score) sammeln."),
    ("Verbündete finden", "Christopher als Begleiter gewinnen."),
    ("Das Geheimnis lösen", "Ins Geheimlabor vordringen und die Wahrheit über die Seuche aufdecken."),
]
runs = []
for t, d in goals:
    runs.append([("> ", 20, GREEN_OK, True), (t + " — ", 20, AMBER_BRT, True),
                 (d, 18, AMBER, False)])
txt(s, Inches(0.85), Inches(1.7), Inches(10.6), Inches(4.6), runs,
    space_after=14, line_spacing=1.05)
footer(s, 5)

# ============================================================
# SLIDE 6 — Besondere Spielkonzepte
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "> ", "Besondere Spielkonzepte")
feats = [
    ("Echtes Terminal-Feeling", "Amber-Phosphor-Look, Typewriter-Effekt, 5 Farbthemen (Amber, Grün, Weiß, Cyan, Rot)."),
    ("Intelligenter Text-Parser", "Versteht viele Verben & Synonyme, reagiert witzig auf Unsinn ('Iss die Waffe')."),
    ("Dynamische Musik & Sound", "Ambient-Tracks im Erkunden, Kampfmusik im Gefecht, zufällige Zombie-Groans."),
    ("Begleiter-System", "Christopher folgt dir, kämpft mit, kann warten — mit eigenen HP."),
    ("Save- & Load-System", "Spielfortschritt als JSON speichern und fortsetzen."),
    ("Versteckte Inhalte", "Geheimlabor mit Boss, Easter-Eggs und ein spezieller Geheim-Gegner."),
]
cx, cy = Inches(0.55), Inches(1.6)
cw, ch = Inches(5.5), Inches(1.5)
gx, gy = Inches(0.3), Inches(0.18)
for i, (t, d) in enumerate(feats):
    col, row = i % 2, i // 2
    l = cx + col * (cw + gx)
    tp = cy + row * (ch + gy)
    add_box(s, l, tp, cw, ch, fill=BG_PANEL, line=AMBER_DIM)
    txt(s, l + Inches(0.25), tp + Inches(0.15), cw - Inches(0.45), ch - Inches(0.3), [
        [("> ", 17, GREEN_OK, True), (t, 17, AMBER_BRT, True)],
        [(d, 14, AMBER, False)],
    ], space_after=5, line_spacing=1.0)
footer(s, 6)

# ============================================================
# SLIDE 7 — NPCs & Begleiter
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "> ", "NPCs & Begleiter")
add_box(s, Inches(0.55), Inches(1.6), Inches(11.05), Inches(2.15), fill=BG_PANEL,
        line=AMBER_DIM)
add_box(s, Inches(0.55), Inches(1.6), Pt(6), Inches(2.15), fill=BLOOD)
txt(s, Inches(0.9), Inches(1.8), Inches(10.4), Inches(1.8), [
    [("Christopher Thomson — der Begleiter", 22, AMBER_BRT, True)],
    [("32 Jahre alt, ~190 cm. Du triffst ihn unterwegs. Wenn du mit ihm sprichst "
      "und sein Vertrauen gewinnst, schließt er sich dir an.", 17, AMBER, False)],
    [("Er folgt dir ('folge mir'), kämpft an deiner Seite mit, kann auf Befehl warten "
      "('bleib hier') und hat eigene Lebenspunkte (HP: 100).", 17, AMBER, False)],
], space_after=7, line_spacing=1.05)
add_box(s, Inches(0.55), Inches(4.0), Inches(11.05), Inches(2.15), fill=BG_PANEL,
        line=AMBER_DIM)
add_box(s, Inches(0.55), Inches(4.0), Pt(6), Inches(2.15), fill=DANGER)
txt(s, Inches(0.9), Inches(4.2), Inches(10.4), Inches(1.8), [
    [("Gegner-Datenbank", 22, AMBER_BRT, True)],
    [("> Toxoplasma-Zombie · Infizierter Mensch", 17, AMBER, False)],
    [("> Mutierter Labor-Leiter (Boss, 280 HP)", 17, AMBER, False)],
    [("> ein geheimer, waffenimmuner Spezial-Gegner mit eigenem Schwachpunkt",
      17, DANGER, True)],
], space_after=6, line_spacing=1.05)
footer(s, 7)

# ============================================================
# SLIDE 8 — Steuerung / Gameplay
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "> ", "Steuerung / Gameplay")
txt(s, Inches(0.85), Inches(1.45), Inches(10.6), Inches(0.6),
    [[("Alles über Texteingabe — kein Maus-Gameplay, reine Kommando-Eingabe:",
       18, AMBER, False)]])
cmds = [
    ("Bewegung", "n / o / s / w / so / nw / hoch / runter"),
    ("Umschauen", "schaue · untersuche · karte · inventar"),
    ("Items", "nimm · lese · esse · öffne · nutze"),
    ("Kampf", "ausrüsten · schieße · schlage · stich · nachladen"),
    ("Begleiter", "spreche · folge mir · bleib hier · begleiter status"),
    ("System", "save · restore · score · hilfe · clear"),
]
add_box(s, Inches(0.55), Inches(2.2), Inches(11.05), Inches(3.9), fill=BG_PANEL,
        line=AMBER_DIM)
runs = []
for t, c in cmds:
    runs.append([(f"{t:<11}", 18, AMBER_BRT, True), ("  ", 18, AMBER, False),
                 (c, 18, GREEN_OK, False)])
txt(s, Inches(0.95), Inches(2.45), Inches(10.3), Inches(3.5), runs,
    space_after=13, line_spacing=1.0)
footer(s, 8)

# ============================================================
# SLIDE 9 — Technischer Aufbau
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "> ", "Technischer Aufbau")
# Modul-Struktur
add_box(s, Inches(0.55), Inches(1.55), Inches(6.4), Inches(4.6), fill=BG_PANEL,
        line=AMBER_DIM)
txt(s, Inches(0.85), Inches(1.75), Inches(5.9), Inches(4.3), [
    [("PROJEKT-STRUKTUR", 19, AMBER_BRT, True)],
    [("dead_world_intro_v_omega.py", 16, GREEN_OK, True), ("  Hauptspiel, ~5500 Z.", 14, AMBER_DIM, False)],
    [("command_handlers.py", 16, GREEN_OK, True), ("  Befehlslogik, ~1900 Z.", 14, AMBER_DIM, False)],
    [("event_handlers.py", 16, GREEN_OK, True), ("  Events & Trigger", 14, AMBER_DIM, False)],
    [("config.py", 16, GREEN_OK, True), ("  Konstanten, Waffen, Gegner", 14, AMBER_DIM, False)],
    [("render_utils.py", 16, GREEN_OK, True), ("  Zeichnen & UI", 14, AMBER_DIM, False)],
    [("Launcher.py", 16, GREEN_OK, True), ("  Start & Updates", 14, AMBER_DIM, False)],
    [("build.bat / PyInstaller", 16, GREEN_OK, True), ("  .exe-Build", 14, AMBER_DIM, False)],
], space_after=8, line_spacing=1.0)
# Technik
add_box(s, Inches(7.15), Inches(1.55), Inches(4.45), Inches(4.6), fill=BG_PANEL,
        line=AMBER_DIM)
txt(s, Inches(7.45), Inches(1.75), Inches(3.95), Inches(4.3), [
    [("TECHNIK", 19, AMBER_BRT, True)],
    bullet_para("Python 3.13", 16),
    bullet_para("Pygame (Grafik & Sound)", 16),
    bullet_para("Eigener Text-Parser", 16),
    bullet_para("JSON Save-System", 16),
    bullet_para("Auflösung bis 1920×1080", 16),
    bullet_para("60 FPS, CRT-Effekt", 16),
    bullet_para("Modularer Aufbau", 16),
    bullet_para("PyInstaller .exe", 16),
], space_after=8, line_spacing=1.0)
footer(s, 9)

# ============================================================
# SLIDE 10 — Schwierigkeiten beim Entwickeln
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "> ", "Schwierigkeiten beim Entwickeln")
diffs = [
    "Eine Hauptdatei wurde riesig (5500+ Zeilen) — Aufteilen in Module war nötig.",
    "Der Text-Parser muss viele Eingaben, Tippfehler & Synonyme abfangen.",
    "~100 Räume konsistent verbinden und Sackgassen vermeiden.",
    "Kampf-Balancing: Waffenschaden vs. Gegner-HP fair einstellen.",
    "Musik & Sounds richtig mischen, ohne dass sie sich überlagern.",
    "Save-/Load-System: alle Spielzustände korrekt sichern und laden.",
    "Bugs finden, die nur in bestimmten Räumen oder Reihenfolgen auftreten.",
]
runs = [bullet_para(d, 18) for d in diffs]
txt(s, Inches(0.85), Inches(1.7), Inches(10.7), Inches(4.5), runs,
    space_after=12, line_spacing=1.05)
footer(s, 10)

# ============================================================
# SLIDE 11 — Was ich gelernt habe
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "> ", "Was ich gelernt habe")
learn = [
    "Wie wichtig Planung & ein Game Design Document vor dem Coden sind.",
    "Code in Module aufzuteilen macht alles übersichtlicher & wartbarer.",
    "Wie man einen eigenen Parser / eine Eingabe-Logik baut.",
    "Probleme im Code systematisch zu finden, zu analysieren und zu beheben.",
    "Wie man Grafik, Sound & Spiellogik mit Pygame verbindet.",
    "Ein fertiges Programm als .exe für andere bereitzustellen.",
    "Dass sich ein Projekt während der Entwicklung stark verändern kann.",
]
runs = [bullet_para(d, 18) for d in learn]
txt(s, Inches(0.85), Inches(1.7), Inches(10.7), Inches(4.5), runs,
    space_after=12, line_spacing=1.05)
footer(s, 11)

# ============================================================
# SLIDE 12 — Zukunft / Verbesserungen
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "> ", "Zukunft / Verbesserungen")
fut = [
    "Mehr Story-Inhalte, Quests und ein richtiges Ende.",
    "Weitere NPCs, Begleiter und einen Handels-/Shop-System.",
    "Crafting-System zum Bauen eigener Items.",
    "Code weiter aufteilen & Dokumentation verbessern.",
    "Mehr automatische Tests, um Bugs früher zu finden.",
    "Save-System robuster machen.",
    "Eventuell ein Tag-/Nacht- oder Hunger-System.",
]
runs = [bullet_para(d, 18) for d in fut]
txt(s, Inches(0.85), Inches(1.7), Inches(10.7), Inches(4.5), runs,
    space_after=12, line_spacing=1.05)
footer(s, 12)

# ============================================================
# SLIDE 13 — Demo
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_box(s, Inches(0.35), Inches(0.35), SW - Inches(0.7), SH - Inches(0.7),
        line=AMBER_DIM, line_w=1.5)
txt(s, 0, Inches(2.0), SW, Inches(0.6),
    [[("[ STARTE DEMO ... ]", 18, AMBER_DIM, False)]], align=PP_ALIGN.CENTER)
txt(s, 0, Inches(2.7), SW, Inches(1.4),
    [[("DEMO DES SPIELS", 64, AMBER_BRT, True)]], align=PP_ALIGN.CENTER,
    anchor=MSO_ANCHOR.MIDDLE)
add_box(s, SW/2 - Inches(2.5), Inches(4.25), Inches(5), Pt(2), fill=BLOOD)
txt(s, 0, Inches(4.5), SW, Inches(0.7),
    [[("Viel Spaß beim Überleben.", 22, AMBER, False)]], align=PP_ALIGN.CENTER)
txt(s, 0, SH - Inches(1.6), SW, Inches(0.6),
    [[("> _", 28, GREEN_OK, True)]], align=PP_ALIGN.CENTER)

out = r"C:\Users\Camilo\Documents\better_game\Dead World.pptx"
prs.save(out)
print("Gespeichert:", out, "| Folien:", len(prs.slides._sldIdLst))
